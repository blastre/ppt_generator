from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import json
import os
from model import ai
from chart_generator import generate_chart, create_chart_directory

# =====================================================
# SLIDE POSITIONING & STYLING CONFIGURATION
# Edit these values to control layout (values in %)
# =====================================================

class SlideConfig:
    """Centralized configuration for slide positioning and styling"""
    
    # SLIDE DIMENSIONS (standard PowerPoint slide)
    SLIDE_WIDTH = 10  # inches
    SLIDE_HEIGHT = 7.5  # inches
    
    # TITLE SLIDE CONFIGURATION
    TITLE_SLIDE = {
        'title': {
            'left_percent': 10,      # % from left edge
            'top_percent': 25,       # % from top edge  
            'width_percent': 80,     # % of slide width
            'height_percent': 15,    # % of slide height
            'font_size': 40,         # Font size in points
            'font_bold': True,
            'alignment': PP_ALIGN.CENTER,
            'color_rgb': (255, 255, 255),  
            'transparency': 0        # 0-100% (0=opaque, 100=transparent)
        },
        'subtitle': {
            'left_percent': 10,
            'top_percent': 60,
            'width_percent': 80,
            'height_percent': 20,
            'font_size': 15,
            'font_bold': False,
            'alignment': PP_ALIGN.LEFT,
            'color_rgb': (190, 180, 180),  
            'transparency': 0
        }
    }
    
    # CONTENT SLIDE CONFIGURATION
    CONTENT_SLIDE = {
        'title': {
            'left_percent': 5,
            'top_percent': 10,
            'width_percent': 90,
            'height_percent': 15,
            'font_size': 24,
            'font_bold': True,
            'alignment': PP_ALIGN.LEFT,
            'color_rgb': (225, 225, 225),  # Blue
            'transparency': 0
        },
        'content': {
            'left_percent': 5,
            'top_percent': 30,
            'width_percent': 90,
            'height_percent': 70,
            'font_size': 14,
            'font_bold': False,
            'alignment': PP_ALIGN.LEFT,
            'color_rgb': (225, 225, 225),    # Dark Gray
            'transparency': 10,
            'margin_left_percent': 2,
            'margin_right_percent': 2,
            'bullet_spacing': 10          # Space after each bullet in points
        }
    }
    
    # CHART SLIDE CONFIGURATION
    CHART_SLIDE = {
        'title': {
            'left_percent': 5,
            'top_percent': 10,
            'width_percent': 90,
            'height_percent': 12,
            'font_size': 24,
            'font_bold': True,
            'alignment': PP_ALIGN.CENTER,
            'color_rgb': (255, 255, 255),  # Blue
            'transparency': 0
        },
        'chart': {
            'left_percent': 70,          # Chart on right half
            'top_percent': 24,
            'width_percent': 45,
            'height_percent': 60,
            'transparency': 0
        },
        'text': {
            'left_percent': 5,           # Text on left half
            'top_percent': 24,
            'width_percent': 40,
            'height_percent': 60,
            'font_size': 12,
            'font_bold': False,
            'alignment': PP_ALIGN.LEFT,
            'color_rgb': (255, 255, 255),   # Dark Gray
            'transparency': 10,
            'margin_left_percent': 2,
            'margin_right_percent': 2,
            'bullet_spacing': 8
        }
    }
    
    # BACKGROUND SETTINGS
    BACKGROUND = {
        'color_rgb': (255, 255, 255),    # White background
        'transparency': 0                # 0-100%
    }

class PPTTemplate:
    """Base class for PowerPoint templates with configurable positioning"""
    
    def __init__(self):
        self.name = "Default"
        self.config = SlideConfig()
        
    def percent_to_inches(self, percent_value, dimension_type='width'):
        """Convert percentage to inches based on slide dimensions"""
        if dimension_type == 'width':
            return Inches((percent_value / 100) * self.config.SLIDE_WIDTH)
        else:  # height
            return Inches((percent_value / 100) * self.config.SLIDE_HEIGHT)
    
    def apply_transparency(self, shape, transparency_percent):
        """Apply transparency to shape (0-100%)"""
        if hasattr(shape, 'fill'):
            shape.fill.transparency = transparency_percent / 100
    
    def apply_title_slide_styling(self, slide, title_text, subtitle_text):
        """Apply styling to title slide with configurable positioning"""
        title_config = self.config.TITLE_SLIDE['title']
        subtitle_config = self.config.TITLE_SLIDE['subtitle']
        
        # Title
        if slide.shapes.title:
            title = slide.shapes.title
            # Position
            title.left = self.percent_to_inches(title_config['left_percent'], 'width')
            title.top = self.percent_to_inches(title_config['top_percent'], 'height')
            title.width = self.percent_to_inches(title_config['width_percent'], 'width')
            title.height = self.percent_to_inches(title_config['height_percent'], 'height')
            
            # Text styling
            title.text = title_text
            title.text_frame.paragraphs[0].font.size = Pt(title_config['font_size'])
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*title_config['color_rgb'])
            title.text_frame.paragraphs[0].font.bold = title_config['font_bold']
            title.text_frame.paragraphs[0].alignment = title_config['alignment']
            
            # Transparency
            self.apply_transparency(title, title_config['transparency'])
        
        # Subtitle
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            # Position
            subtitle.left = self.percent_to_inches(subtitle_config['left_percent'], 'width')
            subtitle.top = self.percent_to_inches(subtitle_config['top_percent'], 'height')
            subtitle.width = self.percent_to_inches(subtitle_config['width_percent'], 'width')
            subtitle.height = self.percent_to_inches(subtitle_config['height_percent'], 'height')
            
            # Text styling
            subtitle.text = subtitle_text
            subtitle.text_frame.paragraphs[0].font.size = Pt(subtitle_config['font_size'])
            subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(*subtitle_config['color_rgb'])
            subtitle.text_frame.paragraphs[0].font.bold = subtitle_config['font_bold']
            subtitle.text_frame.paragraphs[0].alignment = subtitle_config['alignment']
            
            # Transparency
            self.apply_transparency(subtitle, subtitle_config['transparency'])
    
    def apply_content_slide_styling(self, slide, title_text, content_list):
        """Apply styling to content slide with configurable positioning"""
        title_config = self.config.CONTENT_SLIDE['title']
        content_config = self.config.CONTENT_SLIDE['content']
        
        # Title
        if slide.shapes.title:
            title = slide.shapes.title
            # Position
            title.left = self.percent_to_inches(title_config['left_percent'], 'width')
            title.top = self.percent_to_inches(title_config['top_percent'], 'height')
            title.width = self.percent_to_inches(title_config['width_percent'], 'width')
            title.height = self.percent_to_inches(title_config['height_percent'], 'height')
            
            # Text styling
            title.text = title_text
            title.text_frame.paragraphs[0].font.size = Pt(title_config['font_size'])
            title.text_frame.paragraphs[0].font.bold = title_config['font_bold']
            title.text_frame.paragraphs[0].alignment = title_config['alignment']
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*title_config['color_rgb'])
            
            # Transparency
            self.apply_transparency(title, title_config['transparency'])
        
        # Content
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            # Position
            content_placeholder.left = self.percent_to_inches(content_config['left_percent'], 'width')
            content_placeholder.top = self.percent_to_inches(content_config['top_percent'], 'height')
            content_placeholder.width = self.percent_to_inches(content_config['width_percent'], 'width')
            content_placeholder.height = self.percent_to_inches(content_config['height_percent'], 'height')
            
            content_frame = content_placeholder.text_frame
            content_frame.word_wrap = True
            
            # Margins
            content_frame.margin_left = self.percent_to_inches(content_config['margin_left_percent'], 'width')
            content_frame.margin_right = self.percent_to_inches(content_config['margin_right_percent'], 'width')
            
            # Content bullets
            if content_list:
                for i, bullet in enumerate(content_list):
                    if i == 0:
                        content_frame.text = f"• {bullet}"
                    else:
                        p = content_frame.add_paragraph()
                        p.text = f"• {bullet}"
                    
                    paragraph = content_frame.paragraphs[i]
                    paragraph.font.size = Pt(content_config['font_size'])
                    paragraph.font.color.rgb = RGBColor(*content_config['color_rgb'])
                    paragraph.space_after = Pt(content_config['bullet_spacing'])
                    paragraph.alignment = content_config['alignment']
            
            # Transparency
            self.apply_transparency(content_placeholder, content_config['transparency'])
    
    def create_chart_slide(self, prs, title_text, chart_path, bullet_points):
        """Create a chart slide with configurable positioning"""
        title_config = self.config.CHART_SLIDE['title']
        chart_config = self.config.CHART_SLIDE['chart']
        text_config = self.config.CHART_SLIDE['text']
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Title
        title_box = slide.shapes.add_textbox(
            self.percent_to_inches(title_config['left_percent'], 'width'),
            self.percent_to_inches(title_config['top_percent'], 'height'),
            self.percent_to_inches(title_config['width_percent'], 'width'),
            self.percent_to_inches(title_config['height_percent'], 'height')
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(title_config['font_size'])
        title_frame.paragraphs[0].font.bold = title_config['font_bold']
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*title_config['color_rgb'])
        title_frame.paragraphs[0].alignment = title_config['alignment']
        
        # Transparency
        self.apply_transparency(title_box, title_config['transparency'])
        
        # Chart
        if os.path.exists(chart_path):
            chart_shape = slide.shapes.add_picture(
                chart_path,
                self.percent_to_inches(chart_config['left_percent'], 'width'),
                self.percent_to_inches(chart_config['top_percent'], 'height'),
                width=self.percent_to_inches(chart_config['width_percent'], 'width'),
                height=self.percent_to_inches(chart_config['height_percent'], 'height')
            )
            # Transparency
            self.apply_transparency(chart_shape, chart_config['transparency'])
        
        # Bullet points text box
        text_box = slide.shapes.add_textbox(
            self.percent_to_inches(text_config['left_percent'], 'width'),
            self.percent_to_inches(text_config['top_percent'], 'height'),
            self.percent_to_inches(text_config['width_percent'], 'width'),
            self.percent_to_inches(text_config['height_percent'], 'height')
        )
        text_frame = text_box.text_frame
        text_frame.margin_left = self.percent_to_inches(text_config['margin_left_percent'], 'width')
        text_frame.margin_right = self.percent_to_inches(text_config['margin_right_percent'], 'width')
        text_frame.word_wrap = True
        
        if bullet_points:
            for i, bullet in enumerate(bullet_points):
                if i == 0:
                    text_frame.text = f"• {bullet}"
                else:
                    p = text_frame.add_paragraph()
                    p.text = f"• {bullet}"
                
                paragraph = text_frame.paragraphs[i]
                paragraph.font.size = Pt(text_config['font_size'])
                paragraph.font.color.rgb = RGBColor(*text_config['color_rgb'])
                paragraph.space_after = Pt(text_config['bullet_spacing'])
                paragraph.alignment = text_config['alignment']
        
        # Transparency
        self.apply_transparency(text_box, text_config['transparency'])
        
        return slide
    
    def get_contrasting_text_color(self, background_color):
        """Calculate contrasting text color based on background brightness"""
        r, g, b = background_color.r, background_color.g, background_color.b
        luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
        
        if luminance < 0.5:
            return RGBColor(255, 255, 255)  # White
        else:
            return RGBColor(51, 51, 51)     # Dark gray

def get_available_templates():
    """Updated to scan POTX files instead of hardcoded templates"""
    templates = scan_pptx_templates()
    return list(templates.keys())

def create_ppt_skeleton(analysis_result: dict, template_name: str = "default") -> dict:
    """Enhanced skeleton creation with better AI prompting"""
    prompt = f"""Create a professional 5-slide PowerPoint structure for this data analysis:

Question: {analysis_result['original_question']}
Summary: {analysis_result['summary'][:400]}
Template: {template_name}

Requirements:
- Create specific, engaging titles (not generic ones)
- Make bullet points actionable insights, not descriptions
- Focus on business value and key findings
- Keep bullet points to 12-18 words each
- Your response will directly be copy pasted so dont say any extra stuff

Return ONLY this JSON structure:
{{"slides": [
  {{"slide_no": 1, "title": "Executive Summary: [Specific Topic]", "type": "title", "content": ["One compelling insight about the data (one liner only, as short as possible)"]}},
  {{"slide_no": 2, "title": "Data Overview & Methodology", "type": "content", "content": ["Key dataset characteristic", "Analysis approach used", "Important data quality note"]}},
  {{"slide_no": 3, "title": "[Specific Finding Title]", "type": "chart", "content": ["Primary insight discovered", "Supporting evidence", "Business implication"]}},
  {{"slide_no": 4, "title": "[Specific Analysis Title]", "type": "chart", "content": ["Detailed finding", "Trend or pattern identified", "Recommendation"]}},
  {{"slide_no": 5, "title": "Strategic Recommendations", "type": "content", "content": ["Key takeaway", "Recommended action", "Next steps"]}}
]}}"""
    
    response = ai(prompt).strip()
    
    # Enhanced JSON extraction
    if '```' in response:
        response = response.split('```')[1].strip()
        if response.startswith('json'):
            response = response[4:].strip()
    
    start = response.find('{')
    end = response.rfind('}') + 1
    if start != -1 and end > start:
        response = response[start:end]
    
    try:
        return json.loads(response)
    except:
        # Better fallback with specific titles
        return {
            "slides": [
                {"slide_no": 1, "title": f"Executive Summary: {analysis_result['original_question'][:35]}", "type": "title", "content": ["Key insights and findings from comprehensive data analysis"]},
                {"slide_no": 2, "title": "Data Overview & Analysis Approach", "type": "content", "content": ["Dataset scope and key characteristics", "Analytical methodology and tools applied", "Data quality assessment and validation"]},
                {"slide_no": 3, "title": "Primary Data Insights", "type": "chart", "content": ["Most significant patterns discovered in data", "Critical trends affecting business outcomes", "Key performance indicators and metrics"]},
                {"slide_no": 4, "title": "Detailed Analysis Results", "type": "chart", "content": ["In-depth examination of key variables", "Correlation and causation relationships found", "Predictive insights and future implications"]},
                {"slide_no": 5, "title": "Strategic Recommendations", "type": "content", "content": ["Primary business recommendations based on data", "Immediate action items for implementation", "Long-term strategic opportunities identified"]}
            ]
        }

def enhance_slide_content(slide: dict, analysis_result: dict) -> dict:
    """Enhanced content generation with better AI prompting"""
    outline_points = slide.get('content', [])
    
    prompt = f"""Transform these outline points into compelling presentation content:

Slide Title: {slide['title']}
Current Outline: {outline_points}

Context:
- Question: {analysis_result['original_question']}
- Summary: {analysis_result['summary'][:300]}
- Slide Type: {slide['type']}

Requirements:
- Each bullet point: 12-20 words maximum
- Use active voice and strong action verbs
- Include specific data insights where relevant
- Make it business-focused and actionable
- No generic statements

Return exactly {len(outline_points)} enhanced bullet points, one per line, no extra formatting."""
    
    response = ai(prompt).strip()
    lines = response.split('\n')
    
    enhanced_points = []
    for line in lines:
        line = line.strip()
        # Clean formatting
        import re
        line = re.sub(r'^[•\-*►▪▫◦‣\d+\.\s]+', '', line).strip()
        if line and len(line) > 10:
            enhanced_points.append(line)
    
    # Ensure we have content
    slide['content'] = enhanced_points[:4] if enhanced_points else outline_points[:4]
    return slide

def create_presentation(skeleton: dict, analysis_result: dict, output_filename: str, template_name: str = "default") -> str:
    """Updated to use POTX templates and enhanced charts with configurable positioning"""
    
    # Load POTX template and create template instance
    prs = load_pptx_template(template_name)
    template = PPTTemplate()
    chart_dir = create_chart_directory()
    
    print(f"Using template: {template_name}")
    
    for slide_data in skeleton['slides']:
        slide_data = enhance_slide_content(slide_data, analysis_result)
        
        if slide_data['type'] == 'title':
            # Title slide with configurable positioning
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            subtitle_text = '\n'.join(slide_data['content'])
            template.apply_title_slide_styling(slide, slide_data['title'], subtitle_text)
            
        elif slide_data['type'] == 'chart':
            # Enhanced chart slide with configurable positioning
            chart_path = os.path.join(chart_dir, f"chart_{slide_data['slide_no']}.png")
            generate_chart(
                analysis_result['result'], 
                f"{analysis_result['original_question']} - {slide_data['title']}", 
                chart_path
            )
            
            slide = template.create_chart_slide(
                prs, 
                slide_data['title'], 
                chart_path, 
                slide_data['content']
            )
            
        else:
            # Content slide with configurable positioning
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            template.apply_content_slide_styling(slide, slide_data['title'], slide_data['content'])
    
    prs.save(output_filename)
    print(f"Enhanced presentation created with {template_name} template")
    return output_filename
    
def scan_pptx_templates(templates_folder="templates"):
    """Scan for available PPTX template files in templates folder"""
    import glob
    from pathlib import Path
    
    os.makedirs(templates_folder, exist_ok=True)
    template_files = {}
    
    # Look for .pptx files instead of .potx
    pptx_files = glob.glob(os.path.join(templates_folder, "*.pptx"))
    
    for pptx_file in pptx_files:
        template_name = Path(pptx_file).stem.lower().replace(" ", "_").replace("-", "_")
        template_files[template_name] = pptx_file
        
    # Add default if no templates found
    if not template_files:
        template_files["default"] = None
        
    return template_files

def load_pptx_template(template_name="default", templates_folder="templates"):
    """Load PPTX template or create default presentation"""
    available_templates = scan_pptx_templates(templates_folder)
    template_name = template_name.lower()
    
    if template_name in available_templates and available_templates[template_name]:
        try:
            # Load existing .pptx as template
            prs = Presentation(available_templates[template_name])
            # Clear existing slides to use as clean template
            slide_count = len(prs.slides)
            for i in range(slide_count - 1, -1, -1):
                slide_xml = prs.slides._sldIdLst[i]
                prs.part.drop_rel(slide_xml.rId)
                del prs.slides._sldIdLst[i]
            return prs
        except Exception as e:
            print(f"Error loading template {template_name}: {e}")
            return Presentation()
    else:
        return Presentation()
    
def fit_text_to_textbox(text_frame, text, max_font_size=16, min_font_size=10):
    """Intelligently fit text to text frame with proper sizing"""
    text_frame.text = text
    
    # Calculate appropriate font size based on text length and content type
    text_length = len(text)
    line_count = text.count('\n') + 1
    
    if text_length > 300 or line_count > 6:
        font_size = min_font_size
    elif text_length > 200 or line_count > 4:
        font_size = min_font_size + 2
    elif text_length > 100 or line_count > 3:
        font_size = max_font_size - 2
    else:
        font_size = max_font_size
    
    # Apply font size to all paragraphs
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.space_after = Pt(6)
        
    return Pt(font_size)

def create_enhanced_chart_slide(prs, title_text, chart_path, bullet_points, template=None):
    """Legacy function - now uses configurable chart slide creation"""
    if template is None:
        template = PPTTemplate()
    return template.create_chart_slide(prs, title_text, chart_path, bullet_points)